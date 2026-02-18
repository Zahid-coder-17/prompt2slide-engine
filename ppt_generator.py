from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
import os


# ============================================
# 🎨 THEMES
# ============================================

THEMES = {
    "islamic": {
        "bg_color": RGBColor(248, 244, 232),
        "accent_color": RGBColor(0, 90, 60),
        "text_color": RGBColor(40, 40, 40)
    },
    "modern": {
        "bg_color": RGBColor(255, 255, 255),
        "accent_color": RGBColor(30, 120, 200),
        "text_color": RGBColor(30, 30, 30)
    },
    "corporate": {
        "bg_color": RGBColor(245, 247, 250),
        "accent_color": RGBColor(0, 60, 120),
        "text_color": RGBColor(20, 20, 20)
    },
    "dark": {
        "bg_color": RGBColor(25, 25, 25),
        "accent_color": RGBColor(0, 200, 150),
        "text_color": RGBColor(230, 230, 230)
    }
}


# ============================================
# 🎨 BACKGROUND
# ============================================

def set_slide_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ============================================
# 📦 PRESENTATION CREATOR
# ============================================

def create_presentation(slides, theme="modern", output_file="final_presentation.pptx"):

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    theme_data = THEMES.get(theme, THEMES["modern"])

    for slide_data in slides:

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # -------------------------------
        # Background
        # -------------------------------
        set_slide_background(slide, theme_data["bg_color"])

        # -------------------------------
        # Accent Top Bar
        # -------------------------------
        accent_bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0,
            0,
            prs.slide_width,
            Inches(0.35)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = theme_data["accent_color"]
        accent_bar.line.fill.background()

        # -------------------------------
        # Outer Border
        # -------------------------------
        border = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.3),
            Inches(0.6),
            prs.slide_width - Inches(0.6),
            prs.slide_height - Inches(0.9)
        )
        border.fill.background()
        border.line.color.rgb = theme_data["accent_color"]
        border.line.width = Pt(2.5)

        # -------------------------------
        # Title
        # -------------------------------
        title_box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(0.7),
            Inches(11),
            Inches(1)
        )

        tf_title = title_box.text_frame
        tf_title.clear()

        p = tf_title.paragraphs[0]
        run = p.add_run()
        run.text = slide_data["title"]
        run.font.size = Pt(34)
        run.font.bold = True
        run.font.name = "Calibri"
        run.font.color.rgb = theme_data["accent_color"]

        p.alignment = PP_ALIGN.CENTER

        # -------------------------------
        # Bullet Content (LEFT SIDE)
        # -------------------------------
        content_box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.8),
            Inches(6.5),
            Inches(4.8)
        )

        tf = content_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.clear()

        for i, point in enumerate(slide_data["bullet_points"]):

            if i == 0:
                tf.text = point
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = point

            p.level = 0
            p.font.size = Pt(20)
            p.font.name = "Calibri"
            p.font.color.rgb = theme_data["text_color"]

        # -------------------------------
        # Image (RIGHT SIDE)
        # -------------------------------
        image_path = os.path.join(
    "/content/generated_images",
    f"slide_{slide_data['slide_number']}.png"
)


        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path,
                Inches(7.5),
                Inches(1.8),
                width=Inches(4.8),
                height=Inches(4.8)
            )

    prs.save(output_file)
    print(f"\n✅ Presentation saved as {output_file}")



